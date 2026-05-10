"""共享工具：日志、JSON 读写、SHA256 哈希、编码检测、格式转换。"""
import hashlib
import json
import logging
import os
import re
from collections import Counter

_logger_initialized = False


def setup_logging(log_dir):
    global _logger_initialized
    os.makedirs(log_dir, exist_ok=True)
    log_path = os.path.join(log_dir, "pipeline.log")
    logger = logging.getLogger("pipeline")
    if _logger_initialized:
        return logger
    logger.setLevel(logging.DEBUG)
    fh = logging.FileHandler(log_path, encoding="utf-8")
    fh.setLevel(logging.DEBUG)
    fh.setFormatter(logging.Formatter("[%(asctime)s] [%(levelname)s] %(message)s"))
    logger.addHandler(fh)
    ch = logging.StreamHandler()
    ch.setLevel(logging.INFO)
    ch.setFormatter(logging.Formatter("%(message)s"))
    logger.addHandler(ch)
    _logger_initialized = True
    return logger


def load_manifest(output_dir):
    path = os.path.join(output_dir, "manifest.json")
    if os.path.exists(path):
        with open(path, "r", encoding="utf-8") as f:
            return json.load(f)
    return {"files": [], "stats": {}, "duplicate_groups": []}


def save_manifest(manifest, output_dir):
    path = os.path.join(output_dir, "manifest.json")
    with open(path, "w", encoding="utf-8") as f:
        json.dump(manifest, f, ensure_ascii=False, indent=2)


def compute_sha256(filepath, quick=False):
    size = os.path.getsize(filepath)
    h = hashlib.sha256()
    if quick and size > 10 * 1024 * 1024:
        with open(filepath, "rb") as f:
            h.update(f.read(65536))
            f.seek(max(0, size - 65536))
            h.update(f.read(65536))
        h.update(str(size).encode())
    else:
        with open(filepath, "rb") as f:
            for chunk in iter(lambda: f.read(65536), b""):
                h.update(chunk)
    return h.hexdigest()


def detect_encoding(filepath):
    try:
        with open(filepath, "rb") as f:
            raw = f.read(10000)
    except Exception:
        return "utf-8"
    if raw.startswith(b"\xef\xbb\xbf"):
        return "utf-8-sig"
    if raw.startswith(b"\xff\xfe"):
        return "utf-16-le"
    try:
        raw.decode("utf-8")
        return "utf-8"
    except UnicodeDecodeError:
        pass
    try:
        import chardet
        result = chardet.detect(raw)
        return result.get("encoding", "gbk") or "gbk"
    except ImportError:
        return "gbk"


def read_text_file(filepath):
    enc = detect_encoding(filepath)
    with open(filepath, "r", encoding=enc, errors="replace") as f:
        return f.read()


def format_bytes(b):
    if b < 1024:
        return f"{b}B"
    elif b < 1024 * 1024:
        return f"{b/1024:.1f}KB"
    else:
        return f"{b/(1024*1024):.1f}MB"


def format_count(n):
    if n >= 10000:
        return f"{n/10000:.1f}万"
    elif n >= 1000:
        return f"{n/1000:.1f}K"
    return str(n)


def iter_pptx_content(path):
    """共享的 PPTX 幻灯片遍历器，yield (slide_idx, paragraphs, table_rows)。"""
    import pptx
    prs = pptx.Presentation(path)
    for i, slide in enumerate(prs.slides, 1):
        paragraphs = []
        table_rows = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        paragraphs.append(t)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    table_rows.append(cells)
        yield i, paragraphs, table_rows


def iter_ipynb_cells(path):
    """共享的 IPYNB 单元格遍历器，yield (idx, cell_type, source)。"""
    import nbformat
    with open(path, "r", encoding="utf-8") as f:
        nb = nbformat.read(f, as_version=4)
    for i, cell in enumerate(nb.cells):
        yield i, cell.cell_type, cell.source


TAG_TAXONOMY = {
    "course_record": "课程记录册",
    "career_planning": "职业生涯规划",
    "aigc_research": "AIGC研究",
    "course_assignment": "课程作业",
    "group_project": "小组项目",
    "video_production": "视频制作",
    "student_admin": "学生行政",
    "ai_experiment": "AI平台测试",
    "reference_material": "参考材料",
    "creative_writing": "创意写作",
    "math_exercise": "数学练习",
    "python_code": "Python代码",
}

# 个人信息脱敏：先精准替换已知值，再用正则批量替换模式
_ANON_MAP = {
    "陆哲宇轩": "张同学",
    "陆世滨": "陆同学",
    "黄宸宇": "黄同学",
    "韦良健": "韦同学",
    "钟金海": "钟同学",
    "梁权": "梁同学",
    "黄俊杰": "黄同学",
    "251005400222": "2024XXXX01",
    "2515403339": "2024XXXX02",
    "广西财经学院": "XX大学",
    "广西财经": "XX",
    "数科2542": "数科24XX",
    "数科2524": "数科24XX",
    "310宿舍": "XXX宿舍",
}

# 正则脱敏规则 (模式, 替换)
_ANON_PATTERNS = [
    (re.compile(r"1[3-9]\d{9}"), "1XXXXXXXXXX"),
    (re.compile(r"\d{17}[\dXx]"), "XXXXXXXXXXXXXXXXXX"),
    (re.compile(r"[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}"), "***@***.com"),
    (re.compile(r"wxid_[a-zA-Z0-9]+"), "wxid_XXXX"),
    (re.compile(r"\d{4}\.\d{2}\.\d{2}"), "XXXX.XX.XX"),
    (re.compile(r"QQ号\s*\d{5,12}"), "QQ号 XXXXXX"),
    (re.compile(r"QQ\s*\d{5,12}"), "QQ XXXXXX"),
    (re.compile(r"微信号\s*\S+"), "微信号 XXXXX"),
    (re.compile(r"联系电话\s*\d[\d-]*\d"), "联系电话 XXXXXXXXX"),
    (re.compile(r"身份证号\s*\d[\dXx]*"), "身份证号 XXXXXXXXXXXXXXXXXX"),
    (re.compile(r"学\s*号\s*\d{8,12}"), "学 号 2024XXXXXX"),
    (re.compile(r"电子邮箱\s*\S+@\S+"), "电子邮箱 ***@***.com"),
    (re.compile(r"家庭详细地址\s*\S+"), "家庭详细地址 XXXXX"),
    (re.compile(r"籍\s*贯\s*\S{2,10}(?=\s|$)"), "籍 贯 XX省"),
]


def anonymize(text):
    """将文本中的个人信息替换为脱敏占位符。"""
    if not text:
        return text
    # 第一步：精准关键词替换
    for real, fake in _ANON_MAP.items():
        text = text.replace(real, fake)
    # 第二步：正则模式批量替换
    for pattern, replacement in _ANON_PATTERNS:
        text = pattern.sub(replacement, text)
    return text


ALLOWED_EXTENSIONS = {".docx", ".pptx", ".pdf", ".xlsx", ".ipynb", ".csv", ".txt", ".md"}
